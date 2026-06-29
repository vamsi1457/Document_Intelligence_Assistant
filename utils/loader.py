#Extracts text from PDFs, Word docs, CSVs, etc.
import os
import logging
import pandas as pd
import docx
from pptx import Presentation
from langchain_community.document_loaders import PyPDFLoader, TextLoader, CSVLoader
from langchain_core.documents import Document
from .ocr import extract_text_from_image

logger = logging.getLogger(__name__)

def extract_docx_text(file_path: str) -> str:
    """Extracts raw text from a DOCX file."""
    doc = docx.Document(file_path)
    return "\n".join([para.text for para in doc.paragraphs if para.text.strip()])

def extract_pptx_text(file_path: str) -> str:
    """Extracts raw text from all slides in a PPTX file."""
    prs = Presentation(file_path)
    text_runs = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text_runs.append(shape.text)
    return "\n".join(text_runs)

def load_document(file_path: str) -> list[Document]:
    """
    Master router for document ingestion. 
    Reads the file, extracts text, and returns a list of LangChain Documents.
    """
    ext = os.path.splitext(file_path)[1].lower()
    
    try:
        if ext == ".pdf":
            loader = PyPDFLoader(file_path)
            return loader.load()
            
        elif ext in [".txt", ".md"]:
            loader = TextLoader(file_path, encoding="utf-8")
            return loader.load()
            
        elif ext == ".csv":
            loader = CSVLoader(file_path)
            return loader.load()
            
        elif ext in [".docx", ".doc"]:
            text = extract_docx_text(file_path)
            return [Document(page_content=text, metadata={"source": file_path})]
            
        elif ext in [".pptx", ".ppt"]:
            text = extract_pptx_text(file_path)
            return [Document(page_content=text, metadata={"source": file_path})]
            
        elif ext in [".xlsx", ".xls"]:
            # Pandas is significantly more reliable for spreadsheets than raw text extractors
            df = pd.read_excel(file_path)
            text = df.to_string(index=False)
            return [Document(page_content=text, metadata={"source": file_path})]
            
        elif ext in [".png", ".jpg", ".jpeg"]:
            text = extract_text_from_image(file_path)
            return [Document(page_content=text, metadata={"source": file_path})]
            
        else:
            logger.error(f"Unsupported file format attempted: {ext}")
            return []
            
    except Exception as e:
        logger.error(f"Critical failure loading {file_path}. Reason: {e}")
        return []